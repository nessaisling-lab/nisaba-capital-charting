#!/usr/bin/env python
"""Build v13.2 pitch deck as PowerPoint .pptx mirroring the HTML version.

Nisaba Capital Charting · Nisaba Capital Charting
12 slides, ~5 minutes, dark grimoire palette, embedded demo screenshots.

Run: python scripts/build_v13.2_pptx.py
Out: docs/v13.2-pitch-deck.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Palette (matches HTML grimoire deck) ────────────────────────
BG          = RGBColor(0x0F, 0x0A, 0x07)
BG_CARD     = RGBColor(0x1F, 0x1A, 0x14)
GOLD        = RGBColor(0xC9, 0xA4, 0x4A)
GOLD_BRIGHT = RGBColor(0xFF, 0xD8, 0x6B)
GOLD_SOFT   = RGBColor(0x8C, 0x73, 0x33)
INK         = RGBColor(0xF3, 0xEC, 0xD8)
INK_SOFT    = RGBColor(0xC8, 0xB8, 0x8E)
INK_FAINT   = RGBColor(0x8C, 0x7D, 0x57)
GOOD        = RGBColor(0x6F, 0xA8, 0x6B)
BAD         = RGBColor(0xC2, 0x5A, 0x3C)

FONT_DISPLAY = "Georgia"  # serif — closest to Iowan Old Style on Win
FONT_BODY    = "Calibri"

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "screenshots"
OUT = ROOT / "docs" / "v13.2-pitch-deck.pptx"

# ── Slide setup ─────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ─────────────────────────────────────────────────────
def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_multi_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
                   anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """runs = list of dicts with text/size/color/bold/italic/break_before."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for run_spec in runs:
        if run_spec.get("break_before"):
            p = tf.add_paragraph()
            p.alignment = align
            if line_spacing:
                p.line_spacing = line_spacing
        r = p.add_run()
        r.text = run_spec["text"]
        r.font.name = run_spec.get("font", FONT_BODY)
        r.font.size = Pt(run_spec.get("size", 16))
        r.font.bold = run_spec.get("bold", False)
        r.font.italic = run_spec.get("italic", False)
        r.font.color.rgb = run_spec.get("color", INK)
    return tb


def add_border(slide, inset=Inches(0.30), color=GOLD_SOFT, weight_pt=0.75):
    add_rect(
        slide,
        inset, inset,
        SW - 2 * inset, SH - 2 * inset,
        BG, line_color=color, line_width_pt=weight_pt,
    )


def add_slide_marker(slide, num, total, time_label):
    add_text(
        slide, Inches(0.5), Inches(0.18),
        Inches(4), Inches(0.3),
        time_label, size=9, color=INK_FAINT,
        font=FONT_BODY,
    )
    add_text(
        slide, SW - Inches(2.0), Inches(0.18),
        Inches(1.5), Inches(0.3),
        f"{num:02d} / {total:02d}",
        size=9, color=INK_FAINT,
        align=PP_ALIGN.RIGHT, font=FONT_BODY,
    )


def add_eyebrow(slide, x, y, w, text):
    add_text(
        slide, x, y, w, Inches(0.3),
        text.upper(), size=11, bold=True,
        color=GOLD_BRIGHT, font=FONT_BODY,
    )


def add_speaker_cue(slide, text):
    cue_y = SH - Inches(0.95)
    add_rect(
        slide, Inches(0.5), cue_y,
        SW - Inches(1.0), Emu(9525),  # 1px line
        GOLD_SOFT,
    )
    add_text(
        slide, Inches(0.5), cue_y + Inches(0.08),
        SW - Inches(1.0), Inches(0.7),
        f"🎤  {text}",
        size=9, italic=True, color=INK_FAINT,
        font=FONT_BODY,
    )


# ── Total slides constant ───────────────────────────────────────
TOTAL = 12

# ═════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 1, TOTAL, "0:00 — 0:40 · 40s")

add_text(slide, Inches(0.7), Inches(1.2), SW - Inches(1.4), Inches(1.4),
         "Nisaba Capital Charting", size=68, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
add_text(slide, Inches(0.7), Inches(2.6), SW - Inches(1.4), Inches(0.6),
         "A financial intelligence platform at the intersection of "
         "astrology and Wall Street.",
         size=18, italic=True, color=INK_SOFT, align=PP_ALIGN.CENTER)

# Quote block
add_rect(slide, Inches(2.5), Inches(3.5), Emu(38100), Inches(1.2), GOLD)
add_text(
    slide, Inches(2.7), Inches(3.55),
    Inches(8.5), Inches(1.1),
    "“Millionaires don’t use astrology. Billionaires do.”",
    size=28, italic=True, color=INK, font=FONT_DISPLAY,
)
add_text(
    slide, Inches(2.95), Inches(4.55),
    Inches(8.0), Inches(0.4),
    "— attributed to J.P. Morgan",
    size=14, color=INK_SOFT, font=FONT_BODY,
)

# Hedge line
add_multi_text(
    slide, Inches(2.5), Inches(5.0), Inches(8.5), Inches(1.4),
    [
        {"text": "Whether or not he said it, what is documented:  ",
         "size": 13, "color": INK},
        {"text": "Morgan kept the most famous astrologer in America, "
                 "Evangeline Adams, on retainer for investment advice.",
         "size": 13, "color": GOLD_BRIGHT, "bold": True},
        {"text": "", "break_before": True, "size": 12},
        {"text": "This intersection isn’t fringe. It’s hidden. "
                 "I’m Aisling. I believe in astrology. So I built the "
                 "tool these guys had to invent themselves.",
         "size": 12, "italic": True, "color": INK_SOFT, "break_before": True},
    ],
    line_spacing=1.3,
)
add_text(slide, Inches(0.5), SH - Inches(1.3), SW - Inches(1.0), Inches(0.3),
         "AISLING LEIVA  ·  PURSUIT NYC FELLOWSHIP  ·  FINAL PRESENTATION",
         size=10, color=INK_FAINT, align=PP_ALIGN.CENTER)

add_speaker_cue(slide,
    "Open with the quote. Pause 3s. Then the hedge: \"whether or not he said "
    "it, what is documented...\" Land on Adams’ retainer. Drop the closing "
    "tag: \"Not fringe. Hidden. I’m Aisling. I believe. So I built it.\" "
    "Pace 130 wpm — slower than instinct.")


# ═════════════════════════════════════════════════════════════════
# SLIDE 2 — LINEAGE
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 2, TOTAL, "0:25 — 1:00 · 35s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "The Lineage")
add_multi_text(
    slide, Inches(0.7), Inches(1.05), SW - Inches(1.4), Inches(1.3),
    [
        {"text": "This intersection isn’t fringe. It’s ",
         "size": 36, "bold": True, "color": GOLD, "font": FONT_DISPLAY},
        {"text": "hidden.", "size": 36, "bold": True,
         "color": GOLD_BRIGHT, "font": FONT_DISPLAY},
    ],
)

# Verified figures table
rows = [
    ("J.P. Morgan",   "1890s-1913",
     "Kept astrologer Evangeline Adams on retainer. Documented in Adams’ 1926 autobiography."),
    ("W.D. Gann",     "1908",
     "Founded W.D. Gann & Company. Truth of the Stock Tape (1924) praised by The Wall Street Journal."),
    ("Donald Bradley", "1948",
     "Stock Market Prediction. Bradley Siderograph — planetary index for stock turning points. Still tracked today."),
    ("Arch Crawford", "1977-present",
     "Founded Crawford Perspectives. Predicted 1987 crash to the day. #1 market timer 2008+2009 by Hulbert Financial Digest."),
]
table_top = Inches(2.5)
row_h = Inches(0.65)
for i, (name, era, body) in enumerate(rows):
    y = table_top + row_h * i
    # Underline
    add_rect(slide, Inches(0.7), y + row_h - Pt(1),
             SW - Inches(1.4), Emu(9525), GOLD_SOFT)
    add_text(slide, Inches(0.7), y + Inches(0.1),
             Inches(2.2), row_h - Inches(0.1),
             name, size=14, bold=True, color=GOLD_BRIGHT)
    add_text(slide, Inches(2.9), y + Inches(0.1),
             Inches(1.6), row_h - Inches(0.1),
             era, size=12, color=INK_SOFT)
    add_text(slide, Inches(4.5), y + Inches(0.1),
             SW - Inches(5.2), row_h - Inches(0.1),
             body, size=12, color=INK)

# Callout
add_rect(slide, Inches(0.7), Inches(5.5), SW - Inches(1.4), Inches(0.85),
         BG_CARD)
add_rect(slide, Inches(0.7), Inches(5.5), Inches(0.06), Inches(0.85), GOLD)
add_text(
    slide, Inches(0.95), Inches(5.6),
    SW - Inches(1.7), Inches(0.7),
    "Existing tools serving this market are PHP scripts and Python notebooks. "
    "Newsletters charge $1,200/year. The market exists. The production-grade "
    "platform doesn’t.",
    size=14, italic=True, color=INK,
)

add_speaker_cue(slide,
    "Walk 3 names: Morgan, Gann, Crawford. The Crawford line is strongest — "
    "“ranked #1 market timer in 2008 and 2009 by Hulbert Financial Digest.” "
    "Hedge Morgan as “attributed.” Land on the gap: market exists, no platform.")


# ═════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IT IS
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 3, TOTAL, "1:00 — 1:25 · 25s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "What It Is")
add_text(slide, Inches(0.7), Inches(1.05), SW - Inches(1.4), Inches(0.7),
         "Nisaba Capital Charting  ·  three tiers, three audiences, one engine.",
         size=30, bold=True, color=GOLD, font=FONT_DISPLAY)

# Tier strip
tiers = [
    ("FREE",   "For the curious. View any ticker’s natal chart + Lagrange score.", GOLD_SOFT),
    ("PRO",    "For paper traders. Backtest + auto-trading sim.  +37% ROI verified.", GOLD),
    ("MASTER", "For professionals. Sidecar API + OpenBB Workspace integration.", GOLD_SOFT),
]
tier_y = Inches(2.0)
tier_h = Inches(0.85)
tier_w = (SW - Inches(1.4) - Inches(0.4)) / 3
for i, (name, blurb, border_color) in enumerate(tiers):
    x = Inches(0.7) + (tier_w + Inches(0.2)) * i
    add_rect(slide, x, tier_y, tier_w, tier_h, BG_CARD,
             line_color=border_color, line_width_pt=1.0)
    add_text(slide, x + Inches(0.2), tier_y + Inches(0.1),
             tier_w - Inches(0.3), Inches(0.3),
             name, size=11, bold=True, color=GOLD_BRIGHT)
    add_text(slide, x + Inches(0.2), tier_y + Inches(0.4),
             tier_w - Inches(0.3), tier_h - Inches(0.4),
             blurb, size=12, color=INK)

# Three columns: astrology / data / GUI
col_y = Inches(3.2)
col_h = Inches(3.4)
col_w = (SW - Inches(1.4) - Inches(0.4)) / 3
columns = [
    ("ASTROLOGY ENGINE",
     "Wave 9 · 9 modules · 65 tests",
     [
        "Solar Returns (Newton search)",
        "Profections (Hellenistic time-lord)",
        "Planetary returns (Saturn/Jupiter/Mars)",
        "Secondary progressions",
        "9 aspects · 7 patterns · 8 fixed stars",
        "Decans · Sabian Symbols · Eclipses",
        "NASA-grade Swiss Ephemeris",
     ]),
    ("DATA PIPELINE",
     "Wave 7 · 10 native providers",
     [
        "World Bank · IMF · ECB SDMX",
        "BLS · EIA · CFTC COT · OFR FSI",
        "Treasury Direct · CoinGecko · SEC EDGAR",
        "+ Tiingo · Alpha Vantage · Finnhub",
        "+ GDELT · Polymarket · 25 RSS feeds",
        "SQLx compile-time SQL",
        "Tiered priority refresh queue",
     ]),
    ("DESKTOP GUI + SIDECAR",
     "Iced 0.14 · wgpu · axum",
     [
        "3D natal wheel via WGSL shader",
        "Universal pill notification system",
        "Lagrange composite (6 factors)",
        "Backtest with cycle-aligned mode",
        "Paper trading engine",
        "OpenBB Workspace integration",
        "11 REST endpoints + 7 widgets",
     ]),
]
for i, (header, what, bullets) in enumerate(columns):
    x = Inches(0.7) + (col_w + Inches(0.2)) * i
    add_rect(slide, x, col_y, col_w, col_h, BG_CARD,
             line_color=GOLD_SOFT, line_width_pt=0.75)
    add_text(slide, x + Inches(0.2), col_y + Inches(0.15),
             col_w - Inches(0.3), Inches(0.3),
             header, size=10, bold=True, color=GOLD_BRIGHT)
    add_text(slide, x + Inches(0.2), col_y + Inches(0.45),
             col_w - Inches(0.3), Inches(0.4),
             what, size=13, color=GOLD)
    bullet_y = col_y + Inches(0.95)
    for b in bullets:
        add_text(slide, x + Inches(0.2), bullet_y,
                 col_w - Inches(0.3), Inches(0.3),
                 f"·  {b}", size=10, color=INK)
        bullet_y += Inches(0.32)

add_speaker_cue(slide,
    "Three tiers — Free for the curious, Pro for paper traders, Master for "
    "professionals. Three columns = three binaries. Read column headers, then "
    "the bold “what” line. Don’t list the bullets aloud.")


# ═════════════════════════════════════════════════════════════════
# SLIDE 4 — UNIVERSE
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 4, TOTAL, "1:25 — 1:50 · 25s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "Walkthrough · Step 1")
add_text(slide, Inches(0.7), Inches(1.05), SW - Inches(1.4), Inches(0.7),
         "Open the dashboard. Universe tab. 2,112 scored tickers.",
         size=26, bold=True, color=GOLD, font=FONT_DISPLAY)
# workflow pill
add_rect(slide, Inches(0.7), Inches(1.85), Inches(2.4), Inches(0.3), GOLD)
add_text(slide, Inches(0.7), Inches(1.85),
         Inches(2.4), Inches(0.3),
         "▶  WORKFLOW DEMO", size=10, bold=True,
         color=BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Layout: annotation 4.0in left, image 7.5in right
ann_x, ann_y = Inches(0.7), Inches(2.4)
img_x = Inches(5.4)

add_multi_text(
    slide, ann_x, ann_y, Inches(4.5), Inches(4.5),
    [
        {"text": "WHAT YOU’RE SEEING", "size": 10, "bold": True,
         "color": GOLD_BRIGHT},
        {"text": "Sector heat map", "size": 14, "bold": True,
         "color": GOLD, "break_before": True},
        {"text": " at top — color-coded by average astro score. "
                 "Green = optimal, orange = unfavorable.", "size": 14},
        {"text": "2,112-ticker scored universe", "size": 14, "bold": True,
         "color": GOLD, "break_before": True},
        {"text": " table sortable by Astro, Lagrange, Concordance, sector.",
         "size": 14},
        {"text": "Conc column", "size": 14, "bold": True, "color": GOLD,
         "break_before": True},
        {"text": " shows how many of the 6 Lagrange factors agree — "
                 "that’s the high-conviction signal.", "size": 14},
        {"text": "", "break_before": True, "size": 4},
        {"text": "Click any ticker → full detail in <200ms.",
         "size": 12, "italic": True, "color": INK_SOFT,
         "break_before": True},
    ],
    line_spacing=1.3,
)
img_path = SHOTS / "universe-explorer.jpg"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), img_x, Inches(2.3),
                             width=Inches(7.5))

add_speaker_cue(slide,
    "Point to the sector heat map first. Then sweep right to the table. "
    "Highlight the Conc column — “when 4+ factors agree, that’s the "
    "high-conviction signal.” Pick CGCT or any optimal-zone ticker.")


# ═════════════════════════════════════════════════════════════════
# SLIDE 5 — OVERVIEW LAGRANGE
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 5, TOTAL, "1:50 — 2:20 · 30s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "Walkthrough · Step 2")
add_text(slide, Inches(0.7), Inches(1.05), SW - Inches(1.4), Inches(0.7),
         "Six factors converge into one Lagrange score.",
         size=26, bold=True, color=GOLD, font=FONT_DISPLAY)
add_rect(slide, Inches(0.7), Inches(1.85), Inches(2.4), Inches(0.3), GOLD)
add_text(slide, Inches(0.7), Inches(1.85), Inches(2.4), Inches(0.3),
         "▶  WORKFLOW DEMO", size=10, bold=True, color=BG,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_multi_text(
    slide, Inches(0.7), Inches(2.4), Inches(4.5), Inches(4.5),
    [
        {"text": "5 GAUGES, TOP-ALIGNED", "size": 10, "bold": True,
         "color": GOLD_BRIGHT},
        {"text": "Crypto F&G  ·  Equities F&G  ·  Technical  ·  "
                 "Astrology  ·  ★ Lagrange composite",
         "size": 13, "color": INK, "break_before": True},
        {"text": "", "break_before": True, "size": 4},
        {"text": "CHART", "size": 10, "bold": True, "color": GOLD_BRIGHT,
         "break_before": True},
        {"text": "Bollinger Bands · SMA-20 · SMA-50 · volume bars · "
                 "auto-detected technical patterns (double tops, support/resistance)",
         "size": 13, "color": INK, "break_before": True},
        {"text": "", "break_before": True, "size": 4},
        {"text": "90-DAY LAGRANGE HISTORY", "size": 10, "bold": True,
         "color": GOLD_BRIGHT, "break_before": True},
        {"text": "Green band = bullish concordance, red = bearish.",
         "size": 13, "color": INK, "break_before": True},
        {"text": "", "break_before": True, "size": 6},
        {"text": "Lagrange = Tech 20% + Sentiment 20% + Short 15% + "
                 "Insider 15% + Astro 15% + DCF 15%",
         "size": 11, "italic": True, "color": INK_SOFT, "break_before": True},
    ],
    line_spacing=1.3,
)
img_path = SHOTS / "overview-lagrange.jpg"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(5.4), Inches(2.3),
                             width=Inches(7.5))

add_speaker_cue(slide,
    "Sweep across the 5 gauges left-to-right. Land on the Lagrange star. "
    "“This is the composite — six independent signals.” Then point to "
    "the 90-day history strip. “When this stays green for weeks, that’s "
    "structural bullishness, not noise.”")


# ═════════════════════════════════════════════════════════════════
# SLIDE 6 — ASTROLOGY HERO
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 6, TOTAL, "2:20 — 3:00 · 40s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "Walkthrough · Step 3")
add_text(slide, Inches(0.7), Inches(1.05), SW - Inches(1.4), Inches(0.7),
         "Astrology tab — where this gets serious.",
         size=26, bold=True, color=GOLD, font=FONT_DISPLAY)
add_rect(slide, Inches(0.7), Inches(1.85), Inches(2.8), Inches(0.3), GOLD)
add_text(slide, Inches(0.7), Inches(1.85), Inches(2.8), Inches(0.3),
         "▶  WORKFLOW DEMO (HERO)", size=10, bold=True, color=BG,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_multi_text(
    slide, Inches(0.7), Inches(2.4), Inches(4.5), Inches(4.6),
    [
        {"text": "YEAR OF VENUS BADGE", "size": 10, "bold": True,
         "color": GOLD_BRIGHT},
        {"text": "AAPL is in its ", "size": 13, "color": INK,
         "break_before": True},
        {"text": "Year of Venus (10th house · Libra)", "size": 13,
         "bold": True, "color": GOLD},
        {"text": " — Hellenistic profection has rotated Venus into time-lord. "
                 "Aspects involving Venus get +50% Lagrange weight all year.",
         "size": 13, "color": INK},
        {"text": "", "break_before": True, "size": 4},
        {"text": "3D NATAL WHEEL", "size": 10, "bold": True,
         "color": GOLD_BRIGHT, "break_before": True},
        {"text": "Real-time WGSL shader. 13 bodies. Aspect lines glow by "
                 "orb-tightness gradient. Hover any planet → tooltip with "
                 "decan, Sabian, critical degree, OOB.",
         "size": 13, "color": INK, "break_before": True},
        {"text": "", "break_before": True, "size": 4},
        {"text": "LIFECYCLE", "size": 10, "bold": True, "color": GOLD_BRIGHT,
         "break_before": True},
        {"text": "Solar Return + planetary returns + progressed Sun. "
                 "All cached — zero per-render compute.",
         "size": 13, "color": INK, "break_before": True},
    ],
    line_spacing=1.3,
)
img_path = SHOTS / "astrology-aspect-hover.jpg"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(5.4), Inches(2.3),
                             width=Inches(7.5))

add_speaker_cue(slide,
    "Money slide. Slow down. “AAPL is in its Year of Venus right now — "
    "that’s not a number I made up. It’s a 2,000-year-old Hellenistic "
    "technique from Antiochus, first century CE.” Then the wheel: "
    "“rendered in real-time on the GPU via custom WGSL shader.”")


# ═════════════════════════════════════════════════════════════════
# SLIDE 7 — NOTIFICATIONS
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 7, TOTAL, "3:00 — 3:25 · 25s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "Walkthrough · Step 4")
add_text(slide, Inches(0.7), Inches(1.05), SW - Inches(1.4), Inches(0.7),
         "Universal pill notifications + drawer history.",
         size=26, bold=True, color=GOLD, font=FONT_DISPLAY)
add_rect(slide, Inches(0.7), Inches(1.85), Inches(2.4), Inches(0.3), GOLD)
add_text(slide, Inches(0.7), Inches(1.85), Inches(2.4), Inches(0.3),
         "▶  WORKFLOW DEMO", size=10, bold=True, color=BG,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_multi_text(
    slide, Inches(0.7), Inches(2.4), Inches(4.5), Inches(4.5),
    [
        {"text": "THREE PILLS, THREE SIGNALS", "size": 10, "bold": True,
         "color": GOLD_BRIGHT},
        {"text": "ACLS → Optimal  ·  MSFT → Optimal  ·  NVDA → Optimal  "
                 "— Lagrange alerts firing simultaneously.",
         "size": 13, "color": INK, "break_before": True},
        {"text": "", "break_before": True, "size": 4},
        {"text": "Layout NEVER reflows.", "size": 14, "bold": True,
         "color": GOLD, "break_before": True},
        {"text": "Pills slot between the tab strip’s right spacer and "
                 "the bell. Push-down bug killed forever.",
         "size": 13, "color": INK},
        {"text": "", "break_before": True, "size": 4},
        {"text": "Click pill → routes + dismisses in one click.",
         "size": 13, "color": INK, "break_before": True},
        {"text": "Click bell → 24h history. Clear all = wipe.",
         "size": 13, "color": INK, "break_before": True},
        {"text": "", "break_before": True, "size": 6},
        {"text": "Bell rocks 2.4s when active count > 0. Custom Canvas "
                 "widget — no Phosphor codepoint guessing.",
         "size": 11, "italic": True, "color": INK_SOFT, "break_before": True},
    ],
    line_spacing=1.3,
)
img_path = SHOTS / "notifications.png"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(5.4), Inches(2.3),
                             width=Inches(7.5))

add_speaker_cue(slide,
    "Quick on this slide. “Three Optimal alerts firing at once — layout "
    "never moves. That’s intentional engineering. Click any pill, it "
    "routes + dismisses. Click the bell, full 24-hour history.” Move on.")


# ═════════════════════════════════════════════════════════════════
# SLIDE 8 — PROOF +37.06% ROI
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 8, TOTAL, "3:25 — 3:55 · 30s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "The Proof")
add_multi_text(
    slide, Inches(0.7), Inches(1.0), SW - Inches(1.4), Inches(0.85),
    [
        {"text": "+37.06%", "size": 44, "bold": True, "color": GOOD,
         "font": FONT_DISPLAY},
        {"text": "  return on $100K virtual capital.  ",
         "size": 30, "bold": True, "color": GOLD, "font": FONT_DISPLAY},
        {"text": "35 trades. Live forward-test.",
         "size": 22, "color": INK_SOFT, "italic": True},
    ],
)

add_multi_text(
    slide, Inches(0.7), Inches(2.4), Inches(4.7), Inches(4.6),
    [
        {"text": "LIVE PAPER TRADING ENGINE", "size": 10, "bold": True,
         "color": GOLD_BRIGHT},
        {"text": "Starts with $100,000. Auto-opens at Lagrange > 65, exits "
                 "at < 35 or 15% trailing stop.",
         "size": 13, "color": INK, "break_before": True},
        {"text": "", "break_before": True, "size": 6},
        {"text": "Total Value: $137,055.11", "size": 18, "bold": True,
         "color": GOOD, "break_before": True},
        {"text": "Total Return: +37.06%", "size": 18, "bold": True,
         "color": GOOD, "break_before": True},
        {"text": "35 trades  ·  last simulation 2026-05-05",
         "size": 11, "color": INK_SOFT, "break_before": True},
        {"text": "", "break_before": True, "size": 6},
        {"text": "Open positions visible in screenshot:", "size": 13,
         "color": INK, "break_before": True},
        {"text": "  · ABSI +50.63%  · BLSH +17.39%  · AMBA +14.47%  · BUUU +26.60%",
         "size": 12, "color": GOOD, "break_before": True},
        {"text": "", "break_before": True, "size": 6},
        {"text": "Astrology is one of six inputs. The Lagrange strategy "
                 "that produced this return uses it actively. Not cherry-"
                 "picked from history — running every day on real data.",
         "size": 12, "italic": True, "color": INK_SOFT, "break_before": True},
    ],
    line_spacing=1.3,
)
img_path = SHOTS / "paper-trail-37pct.jpg"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(5.6), Inches(2.3),
                             width=Inches(7.3))

add_speaker_cue(slide,
    "Kill shot. Slow down. “$100,000 virtual capital. Plus thirty-seven "
    "percent return. Thirty-five trades. The Lagrange strategy uses astrology "
    "as one of six inputs.” Pause. “This isn’t a backtest cherry-picked "
    "from history. It’s a live forward-testing engine.”")


# ═════════════════════════════════════════════════════════════════
# SLIDE 9 — THESIS
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 9, TOTAL, "3:55 — 4:15 · 20s")

add_eyebrow(slide, Inches(0.7), Inches(1.0), Inches(4), "The Thesis")
add_text(slide, Inches(0.7), Inches(1.4), SW - Inches(1.4), Inches(1.5),
         "Don’t believe.  Measure.",
         size=72, bold=True, color=GOLD, font=FONT_DISPLAY)

add_text(
    slide, Inches(0.7), Inches(3.4), SW - Inches(1.4), Inches(1.0),
    "The thesis isn’t “believe in astrology.” The thesis is — financial "
    "astrology has a documented 100-year track record at the highest levels "
    "of Wall Street. Until now, the techniques weren’t reproducible at scale.",
    size=18, color=INK,
)

# Big callout
add_rect(slide, Inches(0.7), Inches(5.0), SW - Inches(1.4), Inches(1.5),
         BG_CARD)
add_rect(slide, Inches(0.7), Inches(5.0), Inches(0.06), Inches(1.5), GOLD)
add_multi_text(
    slide, Inches(0.95), Inches(5.15), SW - Inches(1.65), Inches(1.3),
    [
        {"text": "Nisaba Capital Charting", "size": 18, "bold": True,
         "color": GOLD_BRIGHT},
        {"text": " is the first open implementation of Hellenistic financial "
                 "astrology with engineering rigor: NASA-grade Swiss Eph "
                 "precision, 132 unit tests against AAPL’s natal chart, a "
                 "backtest engine with cycle-aligned mode (Saturn / Jupiter "
                 "return zones).",
         "size": 14, "color": INK},
        {"text": "", "break_before": True, "size": 6},
        {"text": "Don’t take J.P. Morgan’s word for it. Don’t take "
                 "W.D. Gann’s. ",
         "size": 14, "color": INK, "break_before": True},
        {"text": "Run the backtest yourself.",
         "size": 14, "bold": True, "italic": True, "color": GOLD_BRIGHT},
    ],
    line_spacing=1.3,
)

add_speaker_cue(slide,
    "Slow down. Philosophical pivot. “The thesis isn’t believe — it’s "
    "measure.” Pause. Let it land. Then close: “Run the backtest yourself.”")


# ═════════════════════════════════════════════════════════════════
# SLIDE 10 — NUMBERS
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 10, TOTAL, "4:15 — 4:35 · 20s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "Tech Depth")
add_text(slide, Inches(0.7), Inches(1.05), SW - Inches(1.4), Inches(0.7),
         "30 days  ·  solo  ·  production-grade.",
         size=30, bold=True, color=GOLD, font=FONT_DISPLAY)

stats = [
    ("25k",  "lines of Rust"),
    ("140",  "unit tests"),
    ("10",   "native providers"),
    ("9",    "astrology modules"),
    ("46",   "SQL migrations"),
    ("2",    "GPU shaders"),
    ("3",    "binaries"),
    ("0",    "compiler warnings"),
]
sg_top = Inches(2.2)
sg_h = Inches(2.0)
sg_w = (SW - Inches(1.4) - Inches(0.6)) / 4
for i, (num, label) in enumerate(stats):
    col = i % 4
    row = i // 4
    x = Inches(0.7) + (sg_w + Inches(0.2)) * col
    y = sg_top + (sg_h + Inches(0.2)) * row
    add_rect(slide, x, y, sg_w, sg_h, BG_CARD,
             line_color=GOLD_SOFT, line_width_pt=0.75)
    add_text(slide, x, y + Inches(0.4),
             sg_w, Inches(1.0),
             num, size=44, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
    add_text(slide, x, y + Inches(1.45),
             sg_w, Inches(0.4),
             label.upper(), size=10, color=INK_SOFT,
             align=PP_ALIGN.CENTER)

add_speaker_cue(slide,
    "Read the numbers fast in sequence — they hit harder rapid. End on "
    "“zero compiler warnings.” That’s the engineering credibility punch.")


# ═════════════════════════════════════════════════════════════════
# SLIDE 11 — VISION
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 11, TOTAL, "4:35 — 4:55 · 20s")

add_eyebrow(slide, Inches(0.7), Inches(0.7), Inches(4), "What This Could Be")
add_text(slide, Inches(0.7), Inches(1.05), SW - Inches(1.4), Inches(0.7),
         "Nisaba Capital Charting  —  the platform play.",
         size=30, bold=True, color=GOLD, font=FONT_DISPLAY)

img_path = SHOTS / "charting-capital-landing.png"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(0.7), Inches(2.0),
                             width=Inches(5.5))

add_multi_text(
    slide, Inches(6.7), Inches(2.0), SW - Inches(7.4), Inches(4.5),
    [
        {"text": "Nisaba Capital Charting is the engine. ", "size": 16, "color": INK},
        {"text": "Nisaba Capital Charting", "size": 16, "bold": True,
         "color": GOLD_BRIGHT},
        {"text": " is the brand — a financial intelligence platform at the "
                 "intersection of astrology and Wall Street.",
         "size": 16, "color": INK},
        {"text": "", "break_before": True, "size": 6},
        {"text": "The infrastructure is already production-tier. The "
                 "question isn’t ", "size": 13, "color": INK_SOFT,
         "break_before": True},
        {"text": "can it ship", "size": 13, "italic": True, "color": INK_SOFT},
        {"text": ". The question is ", "size": 13, "color": INK_SOFT},
        {"text": "what audience first", "size": 13, "italic": True,
         "color": INK_SOFT},
        {"text": ":", "size": 13, "color": INK_SOFT},
        {"text": "", "break_before": True, "size": 6},
        {"text": "▸  Retail traders  —  desktop tier ($29/mo)",
         "size": 14, "color": INK, "break_before": True},
        {"text": "▸  Institutional newsletters  —  sidecar API ($1,200/yr "
                 "precedent)", "size": 14, "color": INK, "break_before": True},
        {"text": "▸  Astrology-curious finance professionals  —  the "
                 "W.D. Gann tradition revived",
         "size": 14, "color": INK, "break_before": True},
    ],
    line_spacing=1.4,
)

add_speaker_cue(slide,
    "“Nisaba Capital Charting is the engine. Nisaba Capital Charting is the brand.” Three "
    "audiences = three monetization paths. Don’t dwell — forward-looking, "
    "not prescriptive.")


# ═════════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSE
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_border(slide)
add_slide_marker(slide, 12, TOTAL, "4:55 — 5:00 · 5s")

add_multi_text(
    slide, Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(2.5),
    [
        {"text": "“The stars charted the careers of ",
         "size": 30, "color": INK, "italic": True, "font": FONT_DISPLAY},
        {"text": "Morgan", "size": 30, "bold": True, "italic": True,
         "color": GOLD_BRIGHT, "font": FONT_DISPLAY},
        {"text": ", ", "size": 30, "color": INK, "italic": True,
         "font": FONT_DISPLAY},
        {"text": "Gann", "size": 30, "bold": True, "italic": True,
         "color": GOLD_BRIGHT, "font": FONT_DISPLAY},
        {"text": ", and ", "size": 30, "color": INK, "italic": True,
         "font": FONT_DISPLAY},
        {"text": "Crawford", "size": 30, "bold": True, "italic": True,
         "color": GOLD_BRIGHT, "font": FONT_DISPLAY},
        {"text": ".", "size": 30, "color": INK, "italic": True,
         "font": FONT_DISPLAY},
        {"text": "", "break_before": True, "size": 6},
        {"text": "They’re charting mine too.",
         "size": 30, "italic": True, "color": INK,
         "break_before": True, "font": FONT_DISPLAY},
        {"text": "", "break_before": True, "size": 16},
        {"text": "And I’m building the tool that charts everyone else’s.”",
         "size": 30, "italic": True, "color": INK,
         "break_before": True, "font": FONT_DISPLAY},
    ],
    align=PP_ALIGN.CENTER,
    line_spacing=1.3,
)

add_multi_text(
    slide, Inches(0.7), Inches(5.0), SW - Inches(1.4), Inches(1.4),
    [
        {"text": "Whether you came for the horoscope or stayed for the P&L,",
         "size": 18, "italic": True, "color": INK, "font": FONT_DISPLAY},
        {"text": "the cosmos has been moving markets the whole time.",
         "size": 18, "italic": True, "color": INK,
         "break_before": True, "font": FONT_DISPLAY},
        {"text": "We’re just the first to put it on one screen.",
         "size": 22, "bold": True, "color": GOLD_BRIGHT,
         "break_before": True, "font": FONT_DISPLAY},
    ],
    align=PP_ALIGN.CENTER, line_spacing=1.4,
)

add_text(slide, Inches(0.5), SH - Inches(1.3), SW - Inches(1.0), Inches(0.3),
         "AISLING LEIVA  ·  Nisaba Capital Charting  ·  THANK YOU",
         size=11, color=INK_SOFT, align=PP_ALIGN.CENTER)

add_speaker_cue(slide,
    "Memorize verbatim. Deliver slow. Final beat: “We’re just the first "
    "to put it on one screen.” Pause. Smile. Take questions. Don’t say "
    "“thanks for listening.”")


# ── Save ────────────────────────────────────────────────────────
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"WROTE: {OUT}")
print(f"Slides: {len(prs.slides)}")
